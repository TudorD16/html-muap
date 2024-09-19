import sys
from PyQt5.QtWidgets import (
    QApplication, QMainWindow, QAction, QFileDialog, QTextEdit, QMessageBox,
    QTableWidget, QTableWidgetItem, QVBoxLayout, QWidget, QComboBox, QTabWidget,
    QLabel, QScrollArea
)
from PyQt5.QtGui import QIcon
from docx import Document
import openpyxl
from pptx import Presentation

class OfficeReader(QMainWindow):
    def __init__(self):
        super().__init__()
        self.initUI()

    def initUI(self):
        # Setarea ferestrei principale
        self.setWindowTitle('Office Reader 1.0')
        self.setGeometry(100, 100, 800, 600)

        # Aplicarea temei întunecate
        self.setStyleSheet("""
            QMainWindow {
                background-color: #2e2e2e;
            }
            QTextEdit {
                background-color: #1e1e1e;
                color: #dcdcdc;
            }
            QMenuBar {
                background-color: #1e1e1e;
                color: #dcdcdc;
            }
            QMenuBar::item:selected {
                background-color: #3e3e3e;
            }
            QMenu {
                background-color: #1e1e1e;
                color: #dcdcdc;
            }
            QMenu::item:selected {
                background-color: #3e3e3e;
            }
            QPushButton {
                background-color: #3e3e3e;
                color: #dcdcdc;
                border: 1px solid #555555;
            }
            QPushButton:hover {
                background-color: #4e4e4e;
            }
            QTableWidget {
                background-color: #1e1e1e;
                color: #dcdcdc;
            }
            QTabWidget::pane {
                border: 1px solid #555555;
            }
            QLabel {
                color: #dcdcdc;
            }
            QTextEdit {
                border: none;
            }
        """)

        # Creare zona de text pentru afișarea documentelor Word
        self.textEdit = QTextEdit(self)
        self.textEdit.setReadOnly(True)  # Face documentul Word doar în modul de citire
        self.setCentralWidget(self.textEdit)

        # Meniu
        menubar = self.menuBar()
        fileMenu = menubar.addMenu('File')

        # Adăugare acțiune pentru a deschide fișiere
        openFile = QAction(QIcon('open.png'), 'Open', self)
        openFile.setShortcut('Ctrl+O')
        openFile.triggered.connect(self.showDialog)
        fileMenu.addAction(openFile)

    def showDialog(self):
        # Deschide un dialog pentru a selecta fișiere
        fname, _ = QFileDialog.getOpenFileName(self, 'Open file', '', 
                                               'Word files (*.docx);;Excel files (*.xlsx);;PowerPoint files (*.pptx);;All files (*.*)')
        
        if fname:
            if fname.endswith('.docx'):
                self.openWordFile(fname)
            elif fname.endswith('.xlsx'):
                self.openExcelFile(fname)
            elif fname.endswith('.pptx'):
                self.openPowerPointFile(fname)
            else:
                self.showError("Unsupported file format!")

    def openWordFile(self, filepath):
        try:
            # Restabilește zona de text pentru afișarea documentului Word
            self.textEdit = QTextEdit(self)
            self.textEdit.setReadOnly(True)  # Face documentul Word doar în modul de citire
            self.setCentralWidget(self.textEdit)

            # Citirea fișierului Word
            doc = Document(filepath)
            fullText = [para.text for para in doc.paragraphs]
            self.textEdit.setText('\n'.join(fullText))
        except Exception as e:
            self.showError(f"Error opening Word file: {str(e)}")

    def openExcelFile(self, filepath):
        try:
            # Deschidem fișierul Excel cu data_only=True pentru a obține valorile calculate
            self.wb = openpyxl.load_workbook(filepath, data_only=True)  # Folosim data_only=True pentru a obține valorile calculate
            
            # Creăm un combo box pentru a selecta sheet-ul
            self.sheetSelector = QComboBox(self)
            self.sheetSelector.addItems(self.wb.sheetnames)
            self.sheetSelector.currentIndexChanged.connect(self.loadSheetData)
            
            # Creăm un label pentru a afișa formula
            self.formulaLabel = QLabel(self)
            self.formulaLabel.setStyleSheet("background-color: #2e2e2e; color: #dcdcdc; padding: 5px;")
            
            # Setăm layout-ul și combo box-ul la început
            layout = QVBoxLayout()
            self.tableWidget = QTableWidget()
            self.tableWidget.setEditTriggers(QTableWidget.NoEditTriggers)  # Dezactivează editarea
            self.tableWidget.cellClicked.connect(self.updateFormulaDisplay)
            layout.addWidget(self.sheetSelector)
            layout.addWidget(self.tableWidget)
            layout.addWidget(self.formulaLabel)
            
            container = QWidget()
            container.setLayout(layout)
            self.setCentralWidget(container)
            
            # Încarcă datele pentru sheet-ul selectat inițial
            self.loadSheetData()

        except Exception as e:
            self.showError(f"Error opening Excel file: {str(e)}")

    def loadSheetData(self):
        # Afișează datele din sheet-ul selectat
        selected_sheet = self.sheetSelector.currentText()
        sheet = self.wb[selected_sheet]

        # Setăm numărul de rânduri și coloane
        self.tableWidget.setRowCount(sheet.max_row)
        self.tableWidget.setColumnCount(sheet.max_column)

        # Iterăm prin celulele din sheet și le adăugăm în tabel
        for row_idx, row in enumerate(sheet.iter_rows(values_only=True)):  # Folosim values_only=True pentru a obține doar valorile
            for col_idx, cell in enumerate(row):
                self.tableWidget.setItem(row_idx, col_idx, QTableWidgetItem(str(cell) if cell else ''))


    def updateFormulaDisplay(self, row, column):
        try:
            # Afișează textul complet al celulei selectate
            selected_sheet = self.sheetSelector.currentText()
            sheet = self.wb[selected_sheet]
            cell = sheet.cell(row=row+1, column=column+1)
            
            cell_value = cell.value
            # Verifică dacă celula conține o formulă
            if cell.data_type == 'f':
                cell_value = f"Formula: {cell.value}"
            else:
                cell_value = f"Value: {cell_value}"
            
            self.formulaLabel.setText(cell_value)
        except Exception as e:
            self.formulaLabel.setText("Error displaying cell content")
            print(f"Error updating formula display: {e}")

    def openPowerPointFile(self, filepath):
        try:
            # Creăm un QTabWidget pentru slide-uri
            tabWidget = QTabWidget(self)
            self.setCentralWidget(tabWidget)

            # Deschidem fișierul PowerPoint
            prs = Presentation(filepath)
            
            # Iterăm prin slide-uri și adăugăm fiecare slide în tab-uri
            for slide_idx, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text.strip())
                
                # Creăm un widget scrollabil pentru fiecare slide
                slide_widget = QWidget()
                slide_layout = QVBoxLayout(slide_widget)
                
                # Aplicăm stilul întunecat și textul alb
                label = QLabel("\n".join(slide_text) if slide_text else "No text on this slide")
                label.setWordWrap(True)
                label.setStyleSheet("color: #dcdcdc; background-color: #1e1e1e; padding: 10px;")
                
                scrollArea = QScrollArea()
                scrollArea.setWidget(label)
                scrollArea.setWidgetResizable(True)
                
                slide_layout.addWidget(scrollArea)
                
                # Adăugăm un tab pentru fiecare slide
                tabWidget.addTab(slide_widget, f"Slide {slide_idx + 1}")

        except Exception as e:
            self.showError(f"Error opening PowerPoint file: {str(e)}")


    def showError(self, message):
        msg = QMessageBox()
        msg.setIcon(QMessageBox.Critical)
        msg.setText(message)
        msg.setWindowTitle("Error")
        msg.exec_()

if __name__ == '__main__':
    app = QApplication(sys.argv)
    reader = OfficeReader()
    reader.show()
    sys.exit(app.exec_())
